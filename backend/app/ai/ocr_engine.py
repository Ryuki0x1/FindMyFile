"""
OCR Engine — Extracts text from images using EasyOCR.
Also handles PDF and Office document text extraction natively (no OCR needed for digital PDFs).
GPU-accelerated via PyTorch.
"""

from typing import Optional
from PIL import Image
import numpy as np
import os


class OCREngine:
    """
    Extracts readable text from images and documents.
    
    For images: uses EasyOCR (neural network, handles handwriting + printed text).
    For PDFs: tries native text extraction first (PyMuPDF), falls back to EasyOCR for scanned PDFs.
    For DOCX/PPTX/XLSX: uses python-docx / python-pptx / openpyxl for native text.
    For TXT/MD/CSV: reads directly.
    """

    def __init__(self, languages: list[str] = None):
        self._reader = None
        self._languages = languages or ["en"]

    def load_model(self) -> None:
        """Load EasyOCR reader with improved settings."""
        import easyocr

        self._reader = easyocr.Reader(
            self._languages,
            gpu=True,           # Will fallback to CPU if no CUDA
            model_storage_directory=None,  # use default cache
            download_enabled=True,
        )
        print(f"[OCR] Loaded EasyOCR with languages: {self._languages}")

    def _ensure_loaded(self):
        if self._reader is None:
            self.load_model()

    def extract_text(self, image: Image.Image) -> str:
        """
        Extract all readable text from a PIL Image.
        Uses detail=1 to get confidence scores, filters low-confidence results.
        """
        self._ensure_loaded()

        try:
            # Upscale small images for better OCR accuracy
            w, h = image.size
            if w < 800 or h < 800:
                scale = max(800 / w, 800 / h, 1.5)
                new_w = int(w * scale)
                new_h = int(h * scale)
                image = image.resize((new_w, new_h), Image.Resampling.LANCZOS)

            img_array = np.array(image)

            # detail=1 returns (bbox, text, confidence)
            results = self._reader.readtext(
                img_array,
                detail=1,
                paragraph=False,   # Keep individual words for better keyword matching
                batch_size=8,
                workers=0,
            )

            # Filter by confidence threshold (≥0.4) and join
            texts = [r[1] for r in results if r[2] >= 0.4]
            return " ".join(texts).strip()

        except Exception as e:
            print(f"[OCR] Error: {e}")
            return ""

    def extract_text_from_path(self, filepath: str) -> str:
        """
        Smart text extraction — picks the best method based on file type.
        Order of preference:
          1. Native text extraction (PDF, DOCX, PPTX, XLSX, TXT, MD, CSV)
          2. EasyOCR (images and scanned PDFs)
        """
        ext = os.path.splitext(filepath)[1].lower()

        # --- Native document extraction (fast, accurate, no OCR needed) ---
        if ext == ".pdf":
            text = self._extract_pdf(filepath)
            if text and len(text.strip()) > 20:
                return text
            # Scanned PDF — fall through to EasyOCR

        elif ext in (".docx", ".doc"):
            text = self._extract_docx(filepath)
            if text:
                return text

        elif ext in (".pptx", ".ppt"):
            text = self._extract_pptx(filepath)
            if text:
                return text

        elif ext in (".xlsx", ".xls", ".csv"):
            text = self._extract_xlsx(filepath)
            if text:
                return text

        elif ext in (".txt", ".md", ".rtf"):
            text = self._extract_plaintext(filepath)
            if text:
                return text

        # --- EasyOCR fallback for images and scanned documents ---
        if ext in (".pdf",):
            # Try rendering PDF pages as images and OCR them
            return self._ocr_pdf_pages(filepath)

        # For image files — run EasyOCR directly
        return self._ocr_image_path(filepath)

    # ------------------------------------------------------------------ #
    #  Native extractors                                                   #
    # ------------------------------------------------------------------ #

    def _extract_pdf(self, filepath: str) -> str:
        """Extract text from a digital (non-scanned) PDF using PyMuPDF."""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(filepath)
            texts = []
            for page in doc:
                texts.append(page.get_text("text"))
            doc.close()
            return " ".join(texts).strip()
        except ImportError:
            return ""
        except Exception as e:
            print(f"[OCR] PDF native extract error: {e}")
            return ""

    def _ocr_pdf_pages(self, filepath: str) -> str:
        """Render scanned PDF pages as images and run EasyOCR on them."""
        try:
            import fitz
            self._ensure_loaded()
            doc = fitz.open(filepath)
            texts = []
            for page_num in range(min(len(doc), 5)):  # Max 5 pages for performance
                page = doc[page_num]
                mat = fitz.Matrix(2.0, 2.0)   # 2x zoom for better OCR
                pix = page.get_pixmap(matrix=mat)
                img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
                if pix.n == 4:
                    img_array = img_array[:, :, :3]
                results = self._reader.readtext(img_array, detail=1, paragraph=False)
                page_text = " ".join(r[1] for r in results if r[2] >= 0.4)
                texts.append(page_text)
            doc.close()
            return " ".join(texts).strip()
        except Exception as e:
            print(f"[OCR] PDF page OCR error: {e}")
            return ""

    def _extract_docx(self, filepath: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            import docx
            doc = docx.Document(filepath)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            # Also extract table content
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())
            return " ".join(parts)
        except ImportError:
            return ""
        except Exception as e:
            print(f"[OCR] DOCX extract error: {e}")
            return ""

    def _extract_pptx(self, filepath: str) -> str:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return " ".join(parts)
        except ImportError:
            return ""
        except Exception as e:
            print(f"[OCR] PPTX extract error: {e}")
            return ""

    def _extract_xlsx(self, filepath: str) -> str:
        """Extract text from XLSX/CSV."""
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == ".csv":
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read(5000)
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            parts.append(str(cell).strip())
            wb.close()
            return " ".join(parts)
        except ImportError:
            return ""
        except Exception as e:
            print(f"[OCR] XLSX extract error: {e}")
            return ""

    def _extract_plaintext(self, filepath: str) -> str:
        """Read plain text files directly."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(10000)
        except Exception:
            return ""

    def _ocr_image_path(self, filepath: str) -> str:
        """Run EasyOCR on an image file path."""
        self._ensure_loaded()
        try:
            results = self._reader.readtext(
                filepath,
                detail=1,
                paragraph=False,
                batch_size=8,
                workers=0,
            )
            texts = [r[1] for r in results if r[2] >= 0.4]
            return " ".join(texts).strip()
        except Exception as e:
            print(f"[OCR] Error on {filepath}: {e}")
            return ""
